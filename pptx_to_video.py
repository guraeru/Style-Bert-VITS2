#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX to Video with Voice - Style-Bert-VITS2 (宮舞モカ版)
パワーポイントを音声付き動画に変換する

使い方:
    python pptx_to_video.py <PPTXファイル>
    python pptx_to_video.py <PPTXファイル> -o <出力動画パス>
    python pptx_to_video.py <PPTXファイル> --skip-audio   # テスト用
"""

import os
import re
import sys
import subprocess
import tempfile
import shutil
from pathlib import Path
from typing import Optional, List

import numpy as np
from PIL import Image, ImageDraw, ImageFont

# ===================================================================
# パスの設定（このスクリプト = Style-Bert-VITS2ルート）
# ===================================================================
WORKSPACE = Path(__file__).parent
sys.path.insert(0, str(WORKSPACE))

# PPTX処理
try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"],
                          stdout=subprocess.DEVNULL)
    from pptx import Presentation

# soundfile
try:
    import soundfile as sf
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "soundfile"],
                          stdout=subprocess.DEVNULL)
    import soundfile as sf

# ===================================================================
# 動画設定
# ===================================================================
FRAME_W       = 1280
FRAME_H       = 720
FPS           = 30
PADDING       = 50
BG_COLOR      = (20, 30, 50)        # ダークネイビー
ACCENT_COLOR  = (74, 144, 217)      # 水色
TITLE_COLOR   = (255, 255, 255)     # 白
BODY_COLOR    = (210, 225, 245)     # 薄青白
BAR_COLOR     = (46, 109, 164)      # 青
NUM_COLOR     = (120, 180, 240)     # スライド番号色

MODEL_NAME    = "miyamae_moca"

# ===================================================================
# テキストクリーニング（ナレーション用）
# ===================================================================
_EMOJI_RE = re.compile(
    "["
    "\U0001F300-\U0001F9FF"
    "\U00002600-\U000027BF"
    "\U0000FE00-\U0000FE0F"
    "\U0001F000-\U0001F02F"
    "]+",
    flags=re.UNICODE,
)

_REPLACEMENTS = [
    (r"https?://\S+", ""),          # URL 削除
    (r"[A-Za-z0-9._-]+\.[a-z]{2,}(/\S*)?", ""),  # ドメイン削除
    (r"【(.+?)】", r"\1"),
    (r"[■□◆◇▶▷●○▸]", ""),
        (r"^[\s　]*[・•]\s*", ""),
    (r"[→←↑↓⇒⇐↔]", ""),
    (r"[─━╌╍]+", ""),
    (r"[|｜]", "、"),
    (r"\s{2,}", " "),
]

# ===================================================================
# TTS 向け略語カナ変換
# ===================================================================
def _ascii_token_pat(token: str) -> str:
    """日本語文中でも安全に英数字トークンだけを置換するためのパターン"""
    # \b は日本語の前後で期待どおりに効かないことがあるため lookaround を使う
    return rf"(?<![A-Za-z0-9]){re.escape(token)}(?![A-Za-z0-9])"


_ABBR_MAP = [
    # 固有名詞
    (_ascii_token_pat("VRChat"), "ブイアールチャット"),
    (_ascii_token_pat("Unity"), "ユニティ"),
    (_ascii_token_pat("Blender"), "ブレンダー"),
    (_ascii_token_pat("GitHub"), "ギットハブ"),
    (_ascii_token_pat("Windows"), "ウィンドウズ"),
    (_ascii_token_pat("Mac"), "マック"),
    (_ascii_token_pat("Discord"), "ディスコード"),
    # VRChat 周辺
    (_ascii_token_pat("VRC"), "ブイアールシー"),
    (_ascii_token_pat("VCC"), "ブイシーシー"),
    (_ascii_token_pat("SDK"), "エスディーケー"),
    (_ascii_token_pat("VRCSDK"), "ブイアールシーエスディーケー"),
    (_ascii_token_pat("AV3"), "アバターズスリー"),
    (_ascii_token_pat("Quest"), "クエスト"),
    (_ascii_token_pat("PC"), "ピーシー"),
    (_ascii_token_pat("VR"), "ブイアール"),
    (_ascii_token_pat("FBT"), "フルボディトラッキング"),
    (_ascii_token_pat("FPS"), "エフピーエス"),
    (_ascii_token_pat("STEP"), "ステップ"),
    # 3D/データ周辺
    (_ascii_token_pat("FBX"), "エフビーエックス"),
    (_ascii_token_pat("IK"), "アイケー"),
    (_ascii_token_pat("Rig"), "リグ"),
    (_ascii_token_pat("Mesh"), "メッシュ"),
    (_ascii_token_pat("Bone"), "ボーン"),
    (_ascii_token_pat("Bones"), "ボーン"),
    (_ascii_token_pat("Humanoid"), "ヒューマノイド"),
    (_ascii_token_pat("PhysBone"), "フィジックスボーン"),
    (_ascii_token_pat("PhysBones"), "フィジックスボーン"),
    (_ascii_token_pat("PBone"), "フィジックスボーン"),
    (_ascii_token_pat("PBones"), "フィジックスボーン"),
    # 一般
    (_ascii_token_pat("OK"), "オーケー"),
    (_ascii_token_pat("NG"), "エヌジー"),
    (_ascii_token_pat("Store"), "ストア"),
    (_ascii_token_pat("Booth"), "ブース"),
    (_ascii_token_pat("VRoid"), "ブイロイド"),
    (_ascii_token_pat("VROID"), "ブイロイド"),
    (_ascii_token_pat("Creator Companion"), "クリエイターコンパニオン"),
    (_ascii_token_pat("Avatar Descriptor"), "アバターディスクリプター"),
    (_ascii_token_pat("Eye Movement"), "アイムーブメント"),
    (_ascii_token_pat("GameObject"), "ゲームオブジェクト"),
    (_ascii_token_pat("GameObjects"), "ゲームオブジェクト"),
    (_ascii_token_pat("Read/Write"), "リードライト"),
    (_ascii_token_pat("Read"), "リード"),
    (_ascii_token_pat("Write"), "ライト"),
    (_ascii_token_pat("Publish"), "パブリッシュ"),
    (_ascii_token_pat("Configure"), "コンフィギュア"),
    (_ascii_token_pat("View Position"), "ビューポジション"),
    (_ascii_token_pat("Position"), "ポジション"),
    (_ascii_token_pat("Quality"), "クオリティ"),
    (_ascii_token_pat("Performance"), "パフォーマンス"),
    (_ascii_token_pat("Dynamic Bone"), "ダイナミックボーン"),
    (_ascii_token_pat("Dynamic Bones"), "ダイナミックボーン"),
    (_ascii_token_pat("Dynamic"), "ダイナミック"),
    (_ascii_token_pat("Excellent"), "エクセレント"),
    (_ascii_token_pat("Good"), "グッド"),
    (_ascii_token_pat("Medium"), "ミディアム"),
    (_ascii_token_pat("Very Poor"), "ベリープア"),
    (_ascii_token_pat("Poor"), "プア"),
    (_ascii_token_pat("Build"), "ビルド"),
    (_ascii_token_pat("Import"), "インポート"),
    (_ascii_token_pat("Export"), "エクスポート"),
    (_ascii_token_pat("Asset"), "アセット"),
    (_ascii_token_pat("Assets"), "アセット"),
    (_ascii_token_pat("Plugin"), "プラグイン"),
    (_ascii_token_pat("Plugins"), "プラグイン"),
    (_ascii_token_pat("Component"), "コンポーネント"),
    (_ascii_token_pat("Components"), "コンポーネント"),
    (_ascii_token_pat("Avatar"), "アバター"),
    (_ascii_token_pat("Avatars"), "アバター"),
]


_LETTER_TO_KANA = {
    "A": "エー", "B": "ビー", "C": "シー", "D": "ディー", "E": "イー",
    "F": "エフ", "G": "ジー", "H": "エイチ", "I": "アイ", "J": "ジェー",
    "K": "ケー", "L": "エル", "M": "エム", "N": "エヌ", "O": "オー",
    "P": "ピー", "Q": "キュー", "R": "アール", "S": "エス", "T": "ティー",
    "U": "ユー", "V": "ブイ", "W": "ダブリュー", "X": "エックス", "Y": "ワイ",
    "Z": "ゼット",
}


def _spell_out_acronym(token: str) -> str:
    token = token.upper()
    return "".join(_LETTER_TO_KANA.get(ch, ch) for ch in token)


_ACRONYM_RE = re.compile(r"(?<![A-Za-z0-9])[A-Z]{2,6}(?![A-Za-z0-9])")
_DIGIT_D_RE = re.compile(r"(?<![0-9])([0-9])D(?![A-Za-z0-9])", re.IGNORECASE)
_VERSION_RE = re.compile(r"(?<![A-Za-z0-9])v?(\d+)(?:\.(\d+))+(?![A-Za-z0-9])", re.IGNORECASE)


def convert_abbrs_for_tts(text: str) -> str:
    """英字・略語・表記ゆれをカナ読みに寄せて誤読を減らす"""
    if not text:
        return text

    # 単位表記の最低限の正規化（TTSが詰まりやすい箇所）
    text = re.sub(r"(?<![A-Za-z0-9])(\d+(?:\.\d+)?)m(?![A-Za-z0-9])", r"\1メートル", text)

    # 3D/2D のような表記
    def _digit_d(m):
        d = m.group(1)
        return {"1": "ワン", "2": "ツー", "3": "スリー", "4": "フォー"}.get(d, d) + "ディー"

    text = _DIGIT_D_RE.sub(_digit_d, text)

    # バージョン表記 v1.2.3 -> バージョン1点2点3
    def _version(m):
        s = m.group(0)
        s = re.sub(r"^[Vv]", "", s)
        return "バージョン" + s.replace(".", "点")

    text = _VERSION_RE.sub(_version, text)

    # 既知トークン置換
    for pattern, repl in _ABBR_MAP:
        text = re.sub(pattern, repl, text, flags=re.IGNORECASE)

    # 区切り記号を読みやすい形に寄せる（既知トークン置換の後）
    text = text.replace("&", "アンド")
    text = text.replace("/", "、").replace("\\", "、")
    text = re.sub(r"\s*[-–—]\s*", "、", text)
    text = re.sub(r"\s*\|\s*", "、", text)

    # 未知の英字略語はスペルアウト（例: CPU -> シーピーユー）
    def _acronym(m):
        token = m.group(0)
        # 既知置換で消えていればそのまま
        # （ここでは辞書に無いものだけを汎用変換）
        return _spell_out_acronym(token)

    text = _ACRONYM_RE.sub(_acronym, text)

    # 句読点まわりの余計な空白・連続を整理
    text = re.sub(r"\s*、\s*", "、", text)
    text = re.sub(r"、{2,}", "、", text)
    text = re.sub(r"\s{2,}", " ", text).strip()
    return text

_CIRCLE_NUMS = {
    "①": "1", "②": "2", "③": "3", "④": "4", "⑤": "5",
    "⑥": "6", "⑦": "7", "⑧": "8", "⑨": "9", "⑩": "10",
    "⑪": "11", "⑫": "12",
}

def clean_text(text: str) -> str:
    """ナレーション用にテキストをクリーニング"""
    text = _EMOJI_RE.sub("", text)
    for ch, rep in _CIRCLE_NUMS.items():
        text = text.replace(ch, rep)

    # 範囲表記（1〜4 等）は数字が潰れないように先に処理
    text = re.sub(r"(\d)\s*[〜~]\s*(\d)", r"\1から\2", text)
    # 装飾的な波ダッシュは除去
    text = text.replace("〜", "").replace("~", "")

    for pattern, repl in _REPLACEMENTS:
        text = re.sub(pattern, repl, text)
    text = text.strip()
    return text

def is_table_junk(text: str) -> bool:
    """数値の羅列（表の行）かどうか判定"""
    stripped = re.sub(r"[\d,. 　MBKBGB]+", "", clean_text(text))
    stripped = re.sub(r"[。、\s]+", "", stripped)
    return len(stripped) < 6 and len(text) > 10


# スライドタイトルの繰り返しや単語羅列を除外するためのパターン
_SLIDE_NOISE_PATTERNS = [
    re.compile(r"VRChat\s*3D\s*アバター作成ガイド", re.IGNORECASE),  # 繰り返しタイトル
    re.compile(r"^(情報源|出典)\s*[:：]", re.IGNORECASE),
    re.compile(r"公式ドキュメント\s*\d{4}/\d{1,2}.*更新", re.IGNORECASE),
    re.compile(r"^\d{4}年\d{1,2}月\s*作成$"),
    re.compile(r"^\d{4}/\d{1,2}.*$"),
    re.compile(r"^(Excellent|Good|Medium|Poor|Very\s*Poor)$", re.IGNORECASE),  # ランク単語のみ
    re.compile(r"^\d+$"),                       # 数字のみ
    re.compile(r"^[A-Za-z\s]{1,12}$"),          # 英字のみ・12文字以下（"STEP 1" 等）
    re.compile(r"^STEP\s*\d+$", re.IGNORECASE), # "STEP 1" 形式
]


def _is_noise_line(text: str) -> bool:
    """スライドに大量に出てくる繰り返し・無意味な行を判定"""
    if len(text) < 4:
        return True
    for pat in _SLIDE_NOISE_PATTERNS:
        if pat.search(text):
            return True
    return False

# ===================================================================
# PPTX テキスト抽出
# ===================================================================
def extract_slides(pptx_path: Path) -> List[dict]:
    """各スライドからテキストを抽出する。
    表示用(display_bullets)とナレーション用(narration)は、抽出後にまとめて生成する。
    """
    print(f"[読込] {pptx_path}")
    prs = Presentation(str(pptx_path))
    slides = []

    for i, slide in enumerate(prs.slides):
        title  = ""
        bodies = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            raw = shape.text.strip()
            if not raw:
                continue
            placed = False
            try:
                if shape.is_placeholder:
                    if shape.placeholder_format.type == 1:
                        title  = raw
                        placed = True
            except Exception:
                pass
            if not placed:
                if not title:
                    title = raw
                else:
                    bodies.append(raw)

        slides.append({
            "slide_num": i + 1,
            "title": title,
            "bullets": bodies,
            "display_bullets": [],
            "narration": "",
        })
        print(f"  スライド {i+1:02d}: {title[:42]}")

    # 2nd pass: 表示用とナレーション用を確定
    deck_info = _build_deck_info(slides)
    for idx, sl in enumerate(slides):
        bodies = sl.get("bullets", [])
        sl["display_bullets"] = _select_display_bullets(bodies)
        prev_sl = slides[idx - 1] if idx > 0 else None
        next_sl = slides[idx + 1] if idx + 1 < len(slides) else None
        sl["narration"] = build_narration(
            slide_num=sl["slide_num"],
            title=sl.get("title", ""),
            bodies=bodies,
            prev_slide=prev_sl,
            next_slide=next_sl,
            deck_info=deck_info,
        )

    return slides


def _collect_narration_lines(bodies: List[str]) -> List[str]:
    """ナレーション用に全行をクリーニングして返す"""
    lines = []
    seen = set()
    for b in bodies:
        for raw in b.split("\n"):
            c = clean_text(raw.strip())
            # 先頭の番号だけの列挙を落とす（「1 〜」「2 〜」など）
            c = re.sub(r"^\d+\s+", "", c).strip()
            if not c or is_table_junk(c) or len(c) < 5:
                continue
            if _is_noise_line(c):
                continue
            key = c[:25]
            if key in seen:
                continue
            seen.add(key)
            lines.append(c)
    return lines


def _score_display_line(text: str) -> int:
    score = 0
    ln = len(text)
    if 10 <= ln <= 44:
        score += 6
    elif ln <= 55:
        score += 3
    else:
        score -= 3

    if re.search(r"[ぁ-んァ-ン\u4e00-\u9fff]", text):
        score += 3
    if re.search(r"(手順|準備|設定|確認|注意|ポイント|コツ|やり方|方法|手順|流れ)", text):
        score += 4
    if re.search(r"(する|します|できる|してください|します|しましょう)", text):
        score += 3
    if re.search(r"https?://|\\\\|[A-Za-z0-9]+\.[A-Za-z]{2,}", text):
        score -= 8
    if text.endswith(("：", ":", "-", "—")):
        score -= 2
    return score


def _select_display_bullets(bodies: List[str]) -> List[str]:
    """スライド表示用: 最大 3 項目・短く要点だけ（表の行・繰り返しタイトルは除外）"""
    candidates = []
    seen = set()
    for bi, b in enumerate(bodies):
        for li, raw in enumerate(b.split("\n")):
            c = clean_text(raw.strip())
            if not c or is_table_junk(c) or len(c) < 6:
                continue
            if _is_noise_line(c):
                continue
            key = c[:28]
            if key in seen:
                continue
            seen.add(key)
            candidates.append((bi, li, c, _score_display_line(c)))

    if not candidates:
        return []

    # スコア上位を取りつつ、表示順はスライド内の出現順に戻す
    candidates.sort(key=lambda x: x[3], reverse=True)
    picked = candidates[:6]  # 同スコア対策で少し多めに候補を確保
    picked.sort(key=lambda x: (x[0], x[1]))

    result = []
    for _, __, c, ___ in picked:
        display = c if len(c) <= 50 else c[:48] + "…"
        result.append(display)
        if len(result) >= 3:
            break
    return result


# -------------------------------------------------------------------
# ナレーション生成
# -------------------------------------------------------------------
# スライド番号ごとに関連語を変えてプレゼンターらしさを演出
_INTROS = [
    "",  # スライド1専用
    "続いて{title}です。",
    "ここでは{title}を押さえます。",
    "次は{title}を見ていきましょう。",
    "{title}のポイントです。",
]

_TROUBLE_INTRO = "このスライドは{title}のつまずき対策です。"


def _concept_preamble(clean_title: str) -> str:
    if re.search(r"パフォーマンスランク", clean_title):
        return "ここは基礎パートです。結論から言うと、軽さの目安を先に掴むと、後の最適化が一気に楽になります。"
    if re.search(r"アバターとは", clean_title):
        return "ここは基礎パートです。結論から言うと、作り方の選択肢を先に知ると、遠回りを避けられます。"
    return "ここは基礎パートです。結論から言うと、言葉と全体像を揃えるのが一番の近道です。"

_CONNECTORS = ["最初に、", "次に、", "最後に、"]


_SLIDE_TYPE_RULES = [
    ("agenda", re.compile(r"(目次|本日の流れ|今日の流れ|全体像|講座の流れ|ゴール|概要)", re.IGNORECASE)),
    ("concept", re.compile(r"(とは|基礎|用語|アプローチ|考え方)", re.IGNORECASE)),
    ("setup", re.compile(r"(準備|環境|インストール|導入|前提|必要なもの|セットアップ)", re.IGNORECASE)),
    ("blender", re.compile(r"(Blender|ブレンダー|モデリング|リグ|ウェイト|ボーン|メッシュ)", re.IGNORECASE)),
    ("unity", re.compile(r"(Unity|ユニティ|VCC|Creator\s*Companion|SDK|Package|パッケージ|プロジェクト|アップロード前|アップロード手順|アップロードまで)", re.IGNORECASE)),
    ("trouble", re.compile(r"(注意|NG|トラブル|エラー|よくある|失敗|対処|解決)", re.IGNORECASE)),
    ("summary", re.compile(r"(まとめ|おわり|最後に|チェック|次のステップ)", re.IGNORECASE)),
]


def _classify_slide(title: str, lines: List[str], slide_num: int, total: int) -> str:
    if slide_num == 1:
        return "cover"
    if slide_num == total:
        # 最終スライドはまとめ扱いに寄せる
        return "summary"

    clean_title = clean_text(title)
    # タイトル優先の判定（誤判定を減らす）
    if re.search(r"(Unity|VCC|SDK)", clean_title, re.IGNORECASE):
        return "unity"
    if re.search(r"アップロード", clean_title):
        return "unity"
    if re.search(r"(目次|本日の流れ|今日の流れ|講座の流れ|全体像)", clean_title):
        return "agenda"
    if re.search(r"(とは|基礎|用語|アプローチ|考え方)", clean_title):
        return "concept"

    hay = " ".join([clean_title] + [clean_text(x) for x in lines[:6]])
    for name, pat in _SLIDE_TYPE_RULES:
        if pat.search(hay):
            return name
    return "other"


def _build_deck_info(slides: List[dict]) -> dict:
    """デッキ全体のざっくり構造（パート数など）を推定"""
    total = len(slides)
    types = []
    for sl in slides:
        t = _classify_slide(sl.get("title", ""), _collect_narration_lines(sl.get("bullets", [])), sl["slide_num"], total)
        types.append(t)

    present = set(types)
    ordered = ["setup", "blender", "unity", "trouble"]
    major = [t for t in ordered if t in present]
    return {"total": total, "types": types, "major_parts": major}


def _as_sentence(text: str) -> str:
    """文末を整えて自然な日本語文にする"""
    # 末尾の不要記号を除去
    text = re.sub(r"[、:：\-\s]+$", "", text).strip()
    if not text:
        return ""
    if text.endswith(("。", "！", "？", "す", "い", "る", "れ", "ろ", "ん")):
        return text if text.endswith("。") else text + "。"
    # 数字・記号のみの短い断片は読み飛ばしやすくする
    if len(text) < 6 and not re.search(r"[ぁ-んァ-ン\u4e00-\u9fff]", text):
        return text + "。"
    # 名詞起終のような場合は述語を追加
    return text + "です。"


def build_narration(
    slide_num: int,
    title: str,
    bodies: List[str],
    prev_slide: Optional[dict] = None,
    next_slide: Optional[dict] = None,
    deck_info: Optional[dict] = None,
) -> str:
    """
    スライド内容から、プレゼンターが解説するような自然な日本語ナレーション生成。
    - 画面の文字を丸読みしない
    - スライド番号でイントロ文を変える
    - 重要な点を結合詞でつなぐ
    """
    clean_title = clean_text(title)
    all_lines = _collect_narration_lines(bodies)

    total = deck_info.get("total") if isinstance(deck_info, dict) else None
    slide_type = _classify_slide(clean_title, all_lines, slide_num, total or max(slide_num, 1))

    # スライド1（表紙）: 講座としての宣言を明確化
    if slide_type == "cover":
        t = clean_title or "VRChat 3Dアバター作成"
        major = deck_info.get("major_parts") if isinstance(deck_info, dict) else []
        part_str = ""
        if major:
            jp = {"setup": "準備", "blender": "モデル調整", "unity": "Unityでアップロード", "trouble": "つまずき対策"}
            parts = [(jp.get(x) or x) for x in major[:3]]
            part_str = "今日は、" + "、".join(parts) + "の順で進めます。"
        narration = (
            "ご視聴ありがとうございます。"
            f"この講座では、{t}を、初めての方でも迷わないように一つずつ進めます。"
            + part_str
            + "画面は要点だけに絞っているので、音声の説明に合わせて作業してみてください。"
        )
        return convert_abbrs_for_tts(narration)

    if not clean_title and not all_lines:
        return convert_abbrs_for_tts("次のスライドをご覧ください。")

    # イントロ（タイプに合わせて）
    intro = ""
    if clean_title:
        if slide_type == "trouble":
            intro = _TROUBLE_INTRO.format(title=clean_title)
        else:
            idx = (slide_num - 2) % (len(_INTROS) - 1) + 1
            intro = _INTROS[idx].format(title=clean_title)
    else:
        intro = "ここを押さえていきましょう。"

    # 重要点抽出（最大3）
    key_pts = _pick_key_points(all_lines, 3, slide_type=slide_type)

    # 講座テンプレ: 結論(ゴール)→理由→手順/注意→次へのつなぎ
    parts = [intro]

    if slide_type == "agenda":
        agenda = _pick_agenda_points(all_lines, 3)
        if agenda:
            parts.append("今日は大きく" + str(len(agenda)) + "つ。" + "、".join(agenda) + "の順で進めます。")
        else:
            parts.append("最初に全体の流れを確認してから、作業に入ります。")

        narration = "".join(parts)
        if len(narration) > 260:
            narration = _truncate(narration, 260)
        return convert_abbrs_for_tts(narration).strip()

    elif slide_type == "concept":
        parts.append(_concept_preamble(clean_title))

    elif slide_type == "setup":
        parts.append("ゴールは、必要なツールを揃えて、迷わず作業を始められる状態にすることです。")

    elif slide_type == "blender":
        parts.append("ゴールは、見た目と動きの土台を整えて、アップロードで失敗しない状態にすることです。")

    elif slide_type == "unity":
        parts.append("ここはアップロード手順の要です。順番どおりに設定すれば、失敗をかなり減らせます。")

    elif slide_type == "trouble":
        parts.append("先に結論です。困ったら、原因を一つずつ切り分ければ解決できます。")

    elif slide_type == "summary":
        parts.append("ここまでの要点を短く整理します。")

    # 具体ポイント
    if key_pts:
        for pi, pt in enumerate(key_pts):
            conn = _CONNECTORS[pi] if pi < len(_CONNECTORS) else "次に、"
            parts.append(conn + _as_sentence(pt))
    else:
        parts.append("画面の要点だけ押さえれば大丈夫です。")

    # 次へのつなぎ（連発しない：タイプが切り替わるときだけ）
    if isinstance(deck_info, dict) and next_slide is not None and slide_type not in ("summary",):
        types = deck_info.get("types") or []
        try:
            cur_t = types[slide_num - 1]
            next_t = types[slide_num]
        except Exception:
            cur_t, next_t = None, None
        major = {"setup", "blender", "unity", "trouble", "summary"}
        if next_t and cur_t and next_t != cur_t and (cur_t in major) and (next_t in major) and next_t != "summary":
            parts.append("次はパートが切り替わるので、前提を確認しながら進めましょう。")

    narration = "".join(parts)
    if len(narration) > 320:
        narration = _truncate(narration, 320)
    return convert_abbrs_for_tts(narration).strip()


def _score_key_point(text: str, slide_type: str) -> int:
    score = 0
    ln = len(text)
    if 14 <= ln <= 70:
        score += 6
    elif ln <= 95:
        score += 2
    else:
        score -= 4

    if re.search(r"[ぁ-ん]", text):
        score += 4
    if re.search(r"(する|します|できる|してください|しましょう|確認|設定|入れる|入れます|作る|作ります)", text):
        score += 4
    if re.search(r"(注意|NG|エラー|失敗|対処|回避|ハマる|原因|解決)", text):
        score += 3 if slide_type == "trouble" else 1
    if re.search(r"https?://|\\\\|[A-Za-z0-9]+\.[A-Za-z]{2,}", text):
        score -= 10
    if re.search(r"^[\d\s\-:：/]+$", text):
        score -= 10
    return score


def _pick_key_points(lines: List[str], n: int, slide_type: str = "other") -> List[str]:
    """ナレーション向けに重要度の高い行を最大 n 個返す。"""
    scored = []
    seen = set()
    for li, line in enumerate(lines):
        if len(line) < 10 or len(line) > 110:
            continue
        if _is_noise_line(line):
            continue
        key = line[:28]
        if key in seen:
            continue
        seen.add(key)
        scored.append((li, line, _score_key_point(line, slide_type)))
    if not scored:
        return []
    scored.sort(key=lambda x: x[2], reverse=True)
    top = scored[: max(n, 3)]
    top.sort(key=lambda x: x[0])
    return [x[1] for x in top[:n]]


def _pick_agenda_points(lines: List[str], n: int) -> List[str]:
    pts = []
    for line in lines:
        t = clean_text(line)
        if not t or _is_noise_line(t) or len(t) < 6:
            continue
        # 箇条書きとして読みやすい長さに丸める
        t = re.sub(r"[、。].*$", "", t).strip()
        if len(t) > 28:
            t = t[:26] + "…"
        pts.append(t)
        if len(pts) >= n:
            break
    return pts


def _truncate(text: str, maxlen: int) -> str:
    """最大 maxlen 文字で句点を境に切り詰める"""
    if len(text) <= maxlen:
        return text
    for i in range(maxlen, 0, -1):
        if text[i - 1] in "。！？":
            return text[:i]
    return text[:maxlen] + "。"

# ===================================================================
# スライド画像レンダリング
# ===================================================================
_FONTS: Optional[tuple] = None

def _load_fonts():
    """利用できる日本語フォントを読み込む"""
    candidates = [
        ("C:/Windows/Fonts/meiryo.ttc",   48, 32, 20),
        ("C:/Windows/Fonts/msgothic.ttc",  48, 32, 20),
        ("C:/Windows/Fonts/YuGothM.ttc",   48, 32, 20),
    ]
    for path, ts, bs, ss in candidates:
        if Path(path).exists():
            try:
                return (
                    ImageFont.truetype(path, ts),
                    ImageFont.truetype(path, bs),
                    ImageFont.truetype(path, ss),
                )
            except Exception:
                pass
    f = ImageFont.load_default()
    return f, f, f


def render_slide(slide: dict, total: int) -> Image.Image:
    """スライド1枚を PIL Image として描画。
    表示用に display_bullets (最大 4 項目・少ない情報密度) を使用する。
    """
    global _FONTS
    if _FONTS is None:
        _FONTS = _load_fonts()
    title_f, body_f, small_f = _FONTS

    img  = Image.new("RGB", (FRAME_W, FRAME_H), BG_COLOR)
    draw = ImageDraw.Draw(img)

    n = slide["slide_num"]

    # ヘッダーバー
    draw.rectangle([0, 0, FRAME_W, 100], fill=BAR_COLOR)

    # タイトル（ヘッダー内、垂直中央）
    title_text = slide["title"] or ""
    draw.text((PADDING, 22), title_text, fill=TITLE_COLOR, font=title_f)

    # スライド番号（右上）
    num_str = f"{n:02d} / {total:02d}"
    bbox = small_f.getbbox(num_str)
    nw = bbox[2] - bbox[0]
    draw.text((FRAME_W - PADDING - nw, 72), num_str, fill=NUM_COLOR, font=small_f)

    # 区切り線
    draw.line([(PADDING, 110), (FRAME_W - PADDING, 110)], fill=ACCENT_COLOR, width=2)

    # 本文— display_bullets のみ使用（表示用に選抜済み）
    items = slide.get("display_bullets", [])
    n_items = len(items)
    if n_items == 0:
        pass
    else:
        # 項目数に応じて行間を調整
        content_h = FRAME_H - 110 - 50   # 利用可能高さ
        row_h     = min(content_h // max(n_items, 1), 100)  # 1項目の最大高
        MARKER_X  = PADDING + 8
        TEXT_X    = PADDING + 46
        max_text_w = FRAME_W - TEXT_X - PADDING

        y = 120 + (content_h - row_h * n_items) // 4  # 垂直で少し上寄せ
        for item in items:
            # マーカー
            m_bbox = body_f.getbbox("▸")
            m_h    = m_bbox[3] - m_bbox[1]
            draw.text((MARKER_X, y + (row_h - m_h) // 2), "▸",
                      fill=ACCENT_COLOR, font=body_f)
            # テキスト（wrapして最到10行まで）
            lines = _wrap_text(item, body_f, max_text_w)
            txt_total_h = len(lines) * (body_f.size + 6)
            ty = y + (row_h - txt_total_h) // 2
            for ln in lines:
                draw.text((TEXT_X, ty), ln, fill=BODY_COLOR, font=body_f)
                ty += body_f.size + 6
            y += row_h

    # フッター
    draw.rectangle([0, FRAME_H - 36, FRAME_W, FRAME_H], fill=(25, 40, 65))
    footer = "VRChat アバター作成ガイド  |  宮舞モカ (Style-Bert-VITS2)"
    draw.text((PADDING, FRAME_H - 28), footer, fill=(100, 140, 200), font=small_f)

    return img


def _wrap_text(text: str, font, max_width: int) -> List[str]:
    """テキストを max_width px 以内で折り返す（\n は事前に除去済み前提）"""
    lines   = []
    # \n が残っていれば改行扱い
    paragraphs = text.split("\n")
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        current = ""
        for ch in para:
            test = current + ch
            try:
                bbox = font.getbbox(test)
                w = bbox[2] - bbox[0]
            except Exception:
                w = len(test) * 14  # フォールバック
            if w > max_width:
                if current:
                    lines.append(current)
                current = ch
            else:
                current = test
        if current:
            lines.append(current)
    return lines if lines else [text[:60]]

# ===================================================================
# 音声生成 — TTSModel 直接利用（dubbing_automation.py と同じ方式）
# ===================================================================
_tts_model = None

def init_tts_model(device: str = "cuda"):
    """
    TTSModel を直接インスタンス化する。
    dubbing_automation.py の DubbingAutomation.__init__ と同じアプローチ。
    """
    global _tts_model
    if _tts_model is not None:
        return _tts_model

    import torch
    from style_bert_vits2.tts_model import TTSModel

    # CUDAが使えない場合はCPUにフォールバック
    if device == "cuda" and not torch.cuda.is_available():
        print("  [警告] CUDAが使用不可のため CPU にフォールバックします")
        device = "cpu"

    model_dir = WORKSPACE / "model_assets" / MODEL_NAME

    # 最新の .safetensors を取得（更新日時でソート）
    safetensors = sorted(
        model_dir.glob("*.safetensors"),
        key=lambda p: p.stat().st_mtime,
        reverse=True,
    )
    if not safetensors:
        raise FileNotFoundError(f"safetensors が見つかりません: {model_dir}")

    model_path     = safetensors[0]
    config_path    = model_dir / "config.json"
    style_vec_path = model_dir / "style_vectors.npy"

    for p in (config_path, style_vec_path):
        if not p.exists():
            raise FileNotFoundError(f"必要ファイルなし: {p}")

    print(f"  使用モデル: {model_path.name}")
    _tts_model = TTSModel(
        model_path     = model_path,
        config_path    = config_path,
        style_vec_path = style_vec_path,
        device         = device,
    )
    return _tts_model


def generate_speech(narration: str, out_wav: Path) -> bool:
    """
    ナレーションテキストを WAV に変換して保存する。
    audio_generator.py の generate_audio_for_entry と同じ infer 呼び出しを使用。
    """
    narration = narration.strip()
    if not narration:
        return False

    try:
        sr, audio = _tts_model.infer(
            text         = narration,
            language     = "JP",
            speaker_id   = 0,
            style        = "Neutral",
            style_weight = 1.0,
            length       = 1.0,
        )

        # 1次元 float32 に正規化
        if audio.ndim > 1:
            audio = audio.flatten()
        if audio.dtype == np.int16:
            audio = audio.astype(np.float32) / 32768.0
        else:
            audio = np.clip(audio.astype(np.float32), -1.0, 1.0)

        audio_i16 = (audio * 32767).astype(np.int16)
        out_wav.parent.mkdir(parents=True, exist_ok=True)
        sf.write(str(out_wav), audio_i16, sr)
        return True

    except Exception as e:
        print(f"    [エラー] {e}")
        return False


def get_wav_duration(wav_path: Path) -> float:
    """WAV ファイルの長さ（秒）を返す"""
    try:
        data, sr = sf.read(str(wav_path))
        return len(data) / sr
    except Exception:
        return 5.0

# ===================================================================
# 動画作成 + 音声合成（ffmpeg 使用）
# ===================================================================

def _find_ffmpeg() -> str:
    """ffmpeg 実行ファイルを探す"""
    local = WORKSPACE / "ffmpeg-master-latest-win64-gpl" / "bin" / "ffmpeg.exe"
    if local.exists():
        return str(local)
    found = shutil.which("ffmpeg")
    if found:
        return found
    raise FileNotFoundError(
        "ffmpeg が見つかりません。\n"
        "プロジェクト内 ffmpeg-master-latest-win64-gpl/bin/ を確認してください。"
    )


def _run_ffmpeg(cmd: List[str], label: str = ""):
    """ffmpeg コマンドを実行。失敗時に例外を投げる"""
    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
    )
    if result.returncode != 0:
        raise RuntimeError(
            f"{label} 失敗\n"
            f"CMD: {' '.join(cmd[:5])}...\n"
            f"STDERR: {(result.stderr or '')[-600:]}"
        )
    if label:
        print(f"    ✓ {label}")


def concat_wavs(slides_data, audio_files, fallback_sec: float,
                temp_dir: Path) -> Optional[Path]:
    """
    各スライドの WAV を連結して1本の WAV にまとめる。
    音声がないスライドは無音を挿入する。
    """
    segments = []
    any_ok   = False
    sr_ref   = 44100

    for i, sl in enumerate(slides_data):
        wav = audio_files[i]
        if wav and wav.exists():
            data, sr = sf.read(str(wav), dtype="float32")
            if data.ndim > 1:
                data = data.mean(axis=1)
            sr_ref = sr
            segments.append(data)
            any_ok = True
        else:
            silence = np.zeros(int(fallback_sec * sr_ref), dtype=np.float32)
            segments.append(silence)

    if not any_ok:
        return None

    merged = np.concatenate(segments)
    out    = temp_dir / "merged_audio.wav"
    sf.write(str(out), (merged * 32767).astype(np.int16), sr_ref)
    return out


def create_video(
    slides_data,
    audio_files,
    output_video: Path,
    fallback_sec: float = 6.0,
    temp_dir: Path = None,
):
    """
    スライド画像と音声を組み合わせて MP4 を作成する。

    処理フロー:
      1. スライドごとに PNG を生成
      2. PNG → 無音サブ動画（スライドの尺 = 音声長さ or fallback_sec）
      3. サブ動画を concat
      4. 音声 WAV を連結
      5. 映像 + 音声を ffmpeg で合成
    """
    print("\n[動画組み立て]")
    total   = len(slides_data)
    ffmpeg  = _find_ffmpeg()

    # ── ① スライドごとのサブ動画 ──────────────────────────
    sub_videos = []
    for i, sl in enumerate(slides_data):
        wav = audio_files[i]
        dur = get_wav_duration(wav) if (wav and wav.exists()) else fallback_sec

        # 最短2秒は確保
        dur = max(dur, 2.0)

        # PNG 描画
        img = render_slide(sl, total)
        png = temp_dir / f"frame_{i:03d}.png"
        img.save(str(png))

        sub_mp4 = temp_dir / f"sub_{i:03d}.mp4"
        _run_ffmpeg([
            ffmpeg, "-y",
            "-loop", "1",
            "-framerate", str(FPS),
            "-i", str(png),
            "-t", f"{dur:.3f}",
            "-c:v", "libx264",
            "-pix_fmt", "yuv420p",
            "-preset", "ultrafast",
            str(sub_mp4),
        ], f"スライド {i+1:02d}/{total} 動画化")
        sub_videos.append(sub_mp4)

    # ── ② concat ────────────────────────────────────────────
    concat_txt = temp_dir / "concat.txt"
    concat_txt.write_text(
        "\n".join(f"file '{v.as_posix()}'" for v in sub_videos),
        encoding="utf-8",
    )
    concat_mp4 = temp_dir / "concat.mp4"
    _run_ffmpeg([
        ffmpeg, "-y",
        "-f", "concat", "-safe", "0",
        "-i", str(concat_txt),
        "-c", "copy",
        str(concat_mp4),
    ], "動画を結合")

    # ── ③ 音声 WAV を連結 ────────────────────────────────────
    merged_wav = concat_wavs(slides_data, audio_files, fallback_sec, temp_dir)

    if merged_wav is None:
        # 音声なし → そのまま出力
        shutil.copy(str(concat_mp4), str(output_video))
        print("  (音声なし) 映像のみ出力")
        return True

    # ── ④ 映像 + 音声 合成 ───────────────────────────────────
    _run_ffmpeg([
        ffmpeg, "-y",
        "-i", str(concat_mp4),
        "-i", str(merged_wav),
        "-c:v", "copy",
        "-c:a", "aac",
        "-b:a", "192k",
        "-shortest",
        str(output_video),
    ], "音声を映像に合成")

    print(f"\n  出力先: {output_video}")
    return True

# ===================================================================
# メイン処理
# ===================================================================
def main():
    import argparse

    parser = argparse.ArgumentParser(
        description="PPTX → 音声付き MP4（話者: 宮舞モカ）"
    )
    parser.add_argument("pptx_file",      help="入力 PPTX ファイルパス")
    parser.add_argument("-o", "--output", default=None, help="出力 MP4 パス")
    parser.add_argument("-d", "--duration", type=float, default=6.0,
                        help="音声未生成スライドの表示秒数 (デフォルト: 6)")
    parser.add_argument("--skip-audio",   action="store_true",
                        help="音声生成をスキップ（テスト用）")
    parser.add_argument("--device",       default="cuda",
                        choices=["cpu", "cuda"], help="推論デバイス")
    args = parser.parse_args()

    pptx_path = Path(args.pptx_file)
    if not pptx_path.exists():
        print(f"エラー: ファイルが見つかりません: {pptx_path}")
        return 1

    if args.output:
        output_video = Path(args.output)
    else:
        out_dir = pptx_path.parent / "output_video"
        out_dir.mkdir(exist_ok=True)
        output_video = out_dir / (pptx_path.stem + "_voiced.mp4")
    output_video.parent.mkdir(parents=True, exist_ok=True)

    print(f"\n{'='*62}")
    print(f"  PPTX → 音声付き動画  |  話者: 宮舞モカ")
    print(f"{'='*62}")
    print(f"  入力 : {pptx_path}")
    print(f"  出力 : {output_video}")
    print(f"  デバイス: {args.device}\n")

    # PPTX 解析
    slides_data = extract_slides(pptx_path)
    total = len(slides_data)
    print(f"\n  計 {total} スライドを検出\n")

    # ナレーション確認
    if not args.skip_audio:
        print("[ナレーション一覧]")
        for sl in slides_data:
            print(f"  [{sl['slide_num']:02d}] {sl['narration'][:70]}")
        print()

    # 一時ディレクトリ
    temp_dir = Path(tempfile.mkdtemp(prefix="pptx2video_"))

    try:
        audio_files = []

        # 音声生成
        if not args.skip_audio:
            print("[音声生成]")
            try:
                init_tts_model(args.device)
                print("  ✓ モデルロード完了\n")
            except Exception as e:
                print(f"  [エラー] モデルロード失敗: {e}\n  音声なしで作成します\n")
                args.skip_audio = True

        for i, sl in enumerate(slides_data, 1):
            wav = temp_dir / f"audio_{i:03d}.wav"
            if args.skip_audio:
                audio_files.append(None)
                continue
            narration = sl["narration"]
            print(f"  [{i:02d}/{total}] {narration[:55]}...")
            ok = generate_speech(narration, wav)
            if ok:
                dur = get_wav_duration(wav)
                print(f"         ✓ {wav.name}  ({dur:.1f}秒)")
                audio_files.append(wav)
            else:
                print(f"         ✗ 失敗 → 無音")
                audio_files.append(None)
        print()

        # 動画作成
        create_video(
            slides_data   = slides_data,
            audio_files   = audio_files,
            output_video  = output_video,
            fallback_sec  = args.duration,
            temp_dir      = temp_dir,
        )

        if output_video.exists():
            size_mb = output_video.stat().st_size / 1024 / 1024
            print(f"\n{'='*62}")
            print(f"  ✓ 完成: {output_video}")
            print(f"  サイズ: {size_mb:.1f} MB")
            print(f"{'='*62}\n")
            return 0
        else:
            print("エラー: 出力ファイルが生成されませんでした")
            return 1

    finally:
        try:
            shutil.rmtree(temp_dir)
        except Exception:
            pass


if __name__ == "__main__":
    sys.exit(main())
